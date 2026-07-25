import os
import fitz  # PyMuPDF for PDFs
from docx import Document  # python-docx for Word
from pptx import Presentation  # python-pptx for PowerPoint
import pandas as pd  # pandas for CSV/Excel

class DocumentParser:
    """
    A service class responsible for extracting raw text and metadata 
    from various document formats including PDFs, Word, PPT, Excel, CSV, and TXT.
    """

    @staticmethod
    def parse(file_path: str) -> dict:
        """
        Detects file type by extension and delegates to the appropriate extraction method.
        Returns a dictionary containing raw text and metadata.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found at: {file_path}")
            
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == ".pdf":
            return DocumentParser._parse_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return DocumentParser._parse_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            return DocumentParser._parse_pptx(file_path)
        elif ext in [".xlsx", ".xls"]:
            return DocumentParser._parse_excel(file_path)
        elif ext == ".csv":
            return DocumentParser._parse_csv(file_path)
        elif ext in [".txt", ".md"]:
            return DocumentParser._parse_text(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def _parse_pdf(file_path: str) -> dict:
        text_content = ""
        with fitz.open(file_path) as doc:
            page_count = len(doc)
            for page_num in range(page_count):
                page = doc[page_num]
                text_content += page.get_text()
        return {
            "text": text_content,
            "metadata": {"file_type": "pdf", "page_count": page_count, "source": os.path.basename(file_path)}
        }

    @staticmethod
    def _parse_docx(file_path: str) -> dict:
        doc = Document(file_path)
        text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return {
            "text": text_content,
            "metadata": {"file_type": "docx", "source": os.path.basename(file_path)}
        }

    @staticmethod
    def _parse_pptx(file_path: str) -> dict:
        prs = Presentation(file_path)
        text_content = ""
        slide_count = len(prs.slides)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_content += shape.text + "\n"
        return {
            "text": text_content,
            "metadata": {"file_type": "pptx", "slide_count": slide_count, "source": os.path.basename(file_path)}
        }

    @staticmethod
    def _parse_excel(file_path: str) -> dict:
        df = pd.read_excel(file_path)
        text_content = df.to_string(index=False)
        return {
            "text": text_content,
            "metadata": {"file_type": "excel", "source": os.path.basename(file_path)}
        }

    @staticmethod
    def _parse_csv(file_path: str) -> dict:
        df = pd.read_csv(file_path)
        text_content = df.to_string(index=False)
        return {
            "text": text_content,
            "metadata": {"file_type": "csv", "source": os.path.basename(file_path)}
        }

    @staticmethod
    def _parse_text(file_path: str) -> dict:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text_content = f.read()
        return {
            "text": text_content,
            "metadata": {"file_type": "text", "source": os.path.basename(file_path)}
        }